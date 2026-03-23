"""
Layer 2: Insight overlays — callouts, highlights, badges, brackets, color bands, deltas.
These are added ON TOP of the base data layer (z-order by insertion order in python-pptx).
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from scripts.design_system import (
    DesignSystem, resolve_color,
    CONTENT_LEFT, CONTENT_WIDTH, CONTENT_TOP, CONTENT_HEIGHT,
    FONT_SMALL, FONT_BODY, FONT_FOOTNOTE,
    NAVY, SLATE, WHITE, LIGHT_GREY, TEAL, CRIMSON, MID_GREY,
    POSITIVE_COLOR, NEGATIVE_COLOR,
)
from scripts.utils import set_shape_transparency, set_cell_fill


def apply_overlays(slide, overlays, ds: DesignSystem, context=None):
    """
    Apply a list of overlay specs to a slide.
    context: optional dict with positioning info from the base layer
             (e.g. chart plot area bounds, table cell positions)
    """
    if not overlays:
        return

    ctx = context or {}

    for overlay in overlays:
        overlay_type = overlay.get("type")
        if overlay_type == "callout_annotation":
            _add_callout_annotation(slide, overlay, ds, ctx)
        elif overlay_type == "highlight_box":
            _add_highlight_box(slide, overlay, ds, ctx)
        elif overlay_type == "metric_badge":
            _add_metric_badge(slide, overlay, ds, ctx)
        elif overlay_type == "bracket_group":
            _add_bracket_group(slide, overlay, ds, ctx)
        elif overlay_type == "color_band":
            _add_color_band(slide, overlay, ds, ctx)
        elif overlay_type == "delta_indicator":
            _add_delta_indicator(slide, overlay, ds, ctx)


# ---------------------------------------------------------------------------
# Callout annotation: textbox + optional arrow
# ---------------------------------------------------------------------------
def _add_callout_annotation(slide, spec, ds, ctx):
    """
    Add a callout annotation near a data element.
    spec: {type, text, x, y, width?, height?, color?, arrow_to_x?, arrow_to_y?}
    Or: {type, text, target, position} for chart-relative positioning.
    """
    # Resolve position
    x, y = _resolve_position(spec, ctx)
    width = Inches(spec.get("width", 1.8))
    height = Inches(spec.get("height", 0.45))
    color = resolve_color(spec.get("color", "navy")) if "color" in spec else ds.primary

    # Text box
    txBox = slide.shapes.add_textbox(x, y, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = spec["text"]
    p.font.size = FONT_SMALL
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = ds.font_family
    p.alignment = PP_ALIGN.CENTER

    # Arrow line (if target coordinates provided)
    if "arrow_to_x" in spec and "arrow_to_y" in spec:
        arrow_x = Inches(spec["arrow_to_x"])
        arrow_y = Inches(spec["arrow_to_y"])
        _add_arrow_connector(slide, x + width // 2, y + height, arrow_x, arrow_y, color)


# ---------------------------------------------------------------------------
# Highlight box: semi-transparent rectangle
# ---------------------------------------------------------------------------
def _add_highlight_box(slide, spec, ds, ctx):
    """
    Add a semi-transparent highlight rectangle.
    spec: {type, x, y, width, height, color?, opacity?, label?}
    """
    x = Inches(spec["x"])
    y = Inches(spec["y"])
    width = Inches(spec["width"])
    height = Inches(spec["height"])
    color = resolve_color(spec.get("color", "teal")) if "color" in spec else ds.accent1
    opacity = spec.get("opacity", 20)  # 20% opacity by default

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    set_shape_transparency(shape, opacity)
    shape.line.color.rgb = color
    shape.line.width = Pt(1)

    # Optional label
    if "label" in spec:
        label_box = slide.shapes.add_textbox(x, y - Inches(0.3), width, Inches(0.3))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = spec["label"]
        p.font.size = FONT_FOOTNOTE
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = ds.font_family
        p.alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# Metric badge: big number + label in a box
# ---------------------------------------------------------------------------
def _add_metric_badge(slide, spec, ds, ctx):
    """
    Add a metric callout box.
    spec: {type, value, label, x, y, width?, height?, color?}
    """
    x = Inches(spec.get("x", 10.0))  # Default to top-right area
    y = Inches(spec.get("y", 1.3))
    width = Inches(spec.get("width", 2.0))
    height = Inches(spec.get("height", 0.9))
    color = resolve_color(spec.get("color", "navy")) if "color" in spec else ds.primary

    # Background box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    # Value text (big)
    val_box = slide.shapes.add_textbox(x, y + Inches(0.05), width, Inches(0.5))
    tf = val_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(spec["value"])
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = ds.font_family
    p.alignment = PP_ALIGN.CENTER

    # Label text (small)
    lbl_box = slide.shapes.add_textbox(x, y + Inches(0.48), width, Inches(0.35))
    tf = lbl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = spec["label"]
    p.font.size = FONT_FOOTNOTE
    p.font.color.rgb = LIGHT_GREY
    p.font.name = ds.font_family
    p.alignment = PP_ALIGN.CENTER

    # Optional delta
    if "delta" in spec:
        delta_val = spec["delta"]
        is_pos = not str(delta_val).startswith("-")
        arrow = "\u25B2" if is_pos else "\u25BC"
        delta_color = POSITIVE_COLOR if is_pos else NEGATIVE_COLOR
        delta_box = slide.shapes.add_textbox(x + width - Inches(0.8), y + Inches(0.05), Inches(0.75), Inches(0.3))
        tf = delta_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{arrow} {delta_val}"
        p.font.size = FONT_FOOTNOTE
        p.font.bold = True
        p.font.color.rgb = delta_color
        p.font.name = ds.font_family
        p.alignment = PP_ALIGN.RIGHT


# ---------------------------------------------------------------------------
# Bracket group: line grouping rows/columns with label
# ---------------------------------------------------------------------------
def _add_bracket_group(slide, spec, ds, ctx):
    """
    Add a bracket (vertical line) grouping items with a label.
    spec: {type, x, y_start, y_end, label, color?, side?}
    side: "right" (default) or "left"
    """
    x = Inches(spec["x"])
    y_start = Inches(spec["y_start"])
    y_end = Inches(spec["y_end"])
    color = resolve_color(spec.get("color", "teal")) if "color" in spec else ds.accent1
    bracket_height = y_end - y_start

    # Vertical line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_start, Pt(2.5), bracket_height)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

    # Top cap
    cap_w = Inches(0.15)
    top_cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - cap_w + Pt(2.5), y_start, cap_w, Pt(2))
    top_cap.fill.solid()
    top_cap.fill.fore_color.rgb = color
    top_cap.line.fill.background()

    # Bottom cap
    bot_cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - cap_w + Pt(2.5), y_end - Pt(2), cap_w, Pt(2))
    bot_cap.fill.solid()
    bot_cap.fill.fore_color.rgb = color
    bot_cap.line.fill.background()

    # Label
    label_x = x + Inches(0.15)
    label_y = y_start + (bracket_height - Inches(0.3)) // 2
    label_box = slide.shapes.add_textbox(label_x, label_y, Inches(1.5), Inches(0.3))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = spec["label"]
    p.font.size = FONT_FOOTNOTE
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = ds.font_family


# ---------------------------------------------------------------------------
# Color band: highlight specific table rows
# ---------------------------------------------------------------------------
def _add_color_band(slide, spec, ds, ctx):
    """
    Color-code specific rows in a table.
    spec: {type, row_indices, color, table_ref?}
    ctx must contain "table" key with the table shape.
    """
    table_shape = ctx.get("table")
    if not table_shape:
        return

    table = table_shape.table
    color = resolve_color(spec.get("color", "teal")) if "color" in spec else ds.accent1
    # Make it light — use the color at reduced opacity via a lighter shade
    light_color = _lighten_color(color, 0.85)

    for row_idx in spec.get("row_indices", []):
        if row_idx < len(table.rows):
            for cell in table.rows[row_idx].cells:
                set_cell_fill(cell, light_color)


# ---------------------------------------------------------------------------
# Delta indicator: ▲/▼ with value
# ---------------------------------------------------------------------------
def _add_delta_indicator(slide, spec, ds, ctx):
    """
    Add a standalone delta indicator.
    spec: {type, value, x, y, label?}
    """
    x = Inches(spec["x"])
    y = Inches(spec["y"])
    value = str(spec["value"])
    is_positive = not value.startswith("-")
    arrow = "\u25B2" if is_positive else "\u25BC"
    color = POSITIVE_COLOR if is_positive else NEGATIVE_COLOR

    text = f"{arrow} {value}"
    if "label" in spec:
        text = f"{arrow} {value} {spec['label']}"

    txBox = slide.shapes.add_textbox(x, y, Inches(2), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = FONT_SMALL
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = ds.font_family


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _resolve_position(spec, ctx):
    """Resolve x, y from spec — either explicit or target-based."""
    if "x" in spec and "y" in spec:
        return Inches(spec["x"]), Inches(spec["y"])

    # Target-based positioning (for charts)
    if "target" in spec and "chart_bounds" in ctx:
        bounds = ctx["chart_bounds"]
        categories = ctx.get("categories", [])
        target = spec["target"]
        position = spec.get("position", "above")

        if target in categories:
            idx = categories.index(target)
            num_cats = len(categories)
            plot_left = bounds["left"]
            plot_width = bounds["width"]
            cat_width = plot_width / num_cats
            x = plot_left + cat_width * idx + cat_width * 0.1

            if position == "above":
                y = bounds["top"] - Inches(0.5)
            elif position == "below":
                y = bounds["top"] + bounds["height"] + Inches(0.1)
            else:
                y = bounds["top"] + bounds["height"] / 2
            return Emu(int(x)), Emu(int(y))

    # Fallback
    return CONTENT_LEFT, CONTENT_TOP


def _add_arrow_connector(slide, x1, y1, x2, y2, color):
    """Add a simple line connector between two points."""
    from pptx.oxml.ns import qn
    from lxml import etree

    # Use a freeform connector shape
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, min(x1, x2), min(y1, y2),
        Pt(1), abs(y2 - y1) or Pt(1)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()


def _lighten_color(color, factor=0.85):
    """Lighten an RGBColor by blending with white."""
    r = int(color[0] + (255 - color[0]) * factor)
    g = int(color[1] + (255 - color[1]) * factor)
    b = int(color[2] + (255 - color[2]) * factor)
    return RGBColor(min(r, 255), min(g, 255), min(b, 255))
