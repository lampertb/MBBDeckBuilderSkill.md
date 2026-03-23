"""
Quality validation suite — runs after each slide, prints warnings.
Catches common MBB partner-test failures.
"""

import re
from pptx.util import Inches, Pt, Emu
from scripts.design_system import (
    CONTENT_WIDTH, FONT_HEADLINE, FONT_BODY,
    VIZ_PALETTE, NAVY, SLATE, WHITE, LIGHT_GREY, MID_GREY, TEAL, AMBER, BLUE, CRIMSON, PURPLE, GREEN,
)

# All approved palette colors
APPROVED_COLORS = {
    str(c) for c in [NAVY, SLATE, WHITE, LIGHT_GREY, MID_GREY, TEAL, AMBER, BLUE, CRIMSON, PURPLE, GREEN]
}


def validate_slide(slide, slide_num, slide_data):
    """Run all quality checks on a slide. Returns list of warning strings."""
    warnings = []

    warnings.extend(_check_headline_is_insight(slide_data, slide_num))
    warnings.extend(_check_source_present(slide_data, slide_num))
    warnings.extend(_check_text_overflow(slide, slide_num))
    warnings.extend(_check_shape_count(slide, slide_num, slide_data))

    return warnings


def validate_deck(warnings_by_slide):
    """Print all warnings in a formatted report."""
    total = sum(len(w) for w in warnings_by_slide.values())
    if total == 0:
        print("\n  Quality check: All slides passed.")
        return

    print(f"\n  Quality check: {total} warning(s) found:")
    for slide_num, warnings in sorted(warnings_by_slide.items()):
        for w in warnings:
            print(f"    Slide {slide_num}: {w}")
    print()


def _check_headline_is_insight(slide_data, slide_num):
    """Warn if headline looks like a label instead of a takeaway."""
    warnings = []
    headline = slide_data.get("headline", "")
    slide_type = slide_data.get("type", "")

    # Skip types that don't need insight headlines
    if slide_type in ("title_hero", "section_divider", "qa_page", "busybee_sources", "agenda", "quote_page"):
        return warnings

    if not headline:
        warnings.append("Missing headline — every slide needs a takeaway message")
        return warnings

    words = headline.split()
    if len(words) < 5:
        warnings.append(f'Headline may be a label, not an insight: "{headline}" ({len(words)} words)')

    # Check for verb (simple heuristic)
    has_verb_indicator = any(
        word.endswith(("ed", "ing", "es", "s", "ew", "ove", "ose", "ade"))
        for word in words[1:]  # skip first word
    )
    if not has_verb_indicator and len(words) >= 3:
        # Additional check: common insight verbs
        insight_verbs = {"grew", "declined", "drove", "exceeded", "fell", "rose", "increased",
                        "decreased", "outpaced", "surpassed", "remains", "shifted", "led",
                        "contributed", "accounted", "represents", "shows", "indicates", "suggests",
                        "doubled", "tripled", "halved", "reached", "hit", "missed", "is", "are",
                        "was", "were", "has", "have", "had", "will", "can", "should", "must"}
        if not any(w.lower() in insight_verbs for w in words):
            warnings.append(f'Headline may lack a verb: "{headline}"')

    return warnings


def _check_source_present(slide_data, slide_num):
    """Warn if data-heavy slide has no source."""
    warnings = []
    data_types = {"data_table", "bar_chart", "line_chart", "waterfall_chart", "stacked_bar", "key_stat"}
    if slide_data.get("type") in data_types and not slide_data.get("source"):
        warnings.append("Data slide missing source attribution")
    return warnings


def _check_text_overflow(slide, slide_num):
    """Heuristic overflow check for textboxes."""
    warnings = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = para.text
            if not text:
                continue
            font_size = para.font.size
            if font_size is None:
                font_size = FONT_BODY
            # Rough heuristic: average char width ≈ 0.55 * font_size for proportional fonts
            font_pt = font_size / 12700  # EMU to points
            avg_char_width = font_pt * 0.55
            text_width_pt = len(text) * avg_char_width
            box_width_pt = shape.width / 12700

            if text_width_pt > box_width_pt * 1.1:  # Allow 10% margin for word wrap
                # Only warn for single-line-ish content (no wrapping expected)
                if "\n" not in text and len(text) < 200:
                    pass  # word wrap handles most cases; only flag extreme overflows
                if text_width_pt > box_width_pt * 3.0:
                    warnings.append(f'Possible text overflow: "{text[:50]}..." may not fit')
    return warnings


def _check_shape_count(slide, slide_num, slide_data=None):
    """Warn if slide has too many shapes (too busy)."""
    warnings = []
    count = len(slide.shapes)
    # Complex layouts (matrix, timeline) naturally have more shapes
    slide_type = slide_data.get("type", "") if slide_data else ""
    complex_types = {"two_by_two_matrix", "timeline", "three_column", "agenda",
                     "executive_summary", "key_stat", "two_column"}
    threshold = 40 if slide_type in complex_types else 25
    if count > threshold:
        warnings.append(f"Slide has {count} shapes — may be too busy (target: <{threshold})")
    return warnings
