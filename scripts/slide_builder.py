"""
Core slide builder primitives. Every slide type uses these to place shapes.
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

from scripts.design_system import (
    DesignSystem,
    SLIDE_WIDTH, SLIDE_HEIGHT,
    MARGIN_LEFT, MARGIN_RIGHT, MARGIN_TOP,
    CONTENT_LEFT, CONTENT_WIDTH,
    HEADLINE_TOP, HEADLINE_HEIGHT, HEADLINE_LEFT, HEADLINE_WIDTH,
    DIVIDER_TOP,
    CONTENT_TOP, CONTENT_BOTTOM, CONTENT_HEIGHT,
    FOOTNOTE_TOP, FOOTNOTE_HEIGHT,
    SOURCE_TOP, SOURCE_HEIGHT,
    FONT_HEADLINE, FONT_BODY, FONT_SMALL, FONT_FOOTNOTE, FONT_SOURCE,
    NAVY, SLATE, WHITE, LIGHT_GREY, MID_GREY, TABLE_BORDER,
    TABLE_HEADER_BG, TABLE_HEADER_FG, TABLE_ALT_ROW,
    VIZ_PALETTE, resolve_color,
)
from scripts.utils import set_table_mbb_borders, set_cell_fill, is_numeric


def new_blank_slide(prs):
    """Add and return a blank slide."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)


# ---------------------------------------------------------------------------
# Headline
# ---------------------------------------------------------------------------
def add_headline(slide, text, ds: DesignSystem):
    """Add the takeaway-message headline at top of slide. Auto-shrinks for long text."""
    txBox = slide.shapes.add_textbox(HEADLINE_LEFT, HEADLINE_TOP, HEADLINE_WIDTH, HEADLINE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Enable auto-fit (shrink text to fit box) via XML
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        from lxml import etree
        # Remove existing autofit elements
        for child in list(bodyPr):
            if etree.QName(child.tag).localname in ("noAutofit", "normAutofit", "spAutoFit"):
                bodyPr.remove(child)
        norm = etree.SubElement(bodyPr, qn("a:normAutofit"))
        norm.set("fontScale", "100000")  # allow shrink

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = ds.font_headline
    p.font.bold = True
    p.font.color.rgb = ds.primary
    p.font.name = ds.font_family
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(0)
    return txBox


def add_divider_line(slide, ds: DesignSystem):
    """Add a thin horizontal line below the headline."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        HEADLINE_LEFT,
        DIVIDER_TOP,
        HEADLINE_WIDTH,
        Pt(1.5),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ds.accent1
    line.line.fill.background()
    return line


# ---------------------------------------------------------------------------
# Source & Footnotes
# ---------------------------------------------------------------------------
def add_source(slide, text, ds: DesignSystem):
    """Add source attribution at bottom-left."""
    if not text:
        return None
    txBox = slide.shapes.add_textbox(CONTENT_LEFT, SOURCE_TOP, CONTENT_WIDTH, SOURCE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Source: {text}"
    p.font.size = FONT_SOURCE
    p.font.italic = True
    p.font.color.rgb = MID_GREY
    p.font.name = ds.font_family
    return txBox


def add_footnotes(slide, footnotes, ds: DesignSystem):
    """Add footnote(s) above source line."""
    if not footnotes:
        return None
    text = "  ".join(f"({i+1}) {fn}" for i, fn in enumerate(footnotes))
    txBox = slide.shapes.add_textbox(CONTENT_LEFT, FOOTNOTE_TOP, CONTENT_WIDTH, FOOTNOTE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = FONT_FOOTNOTE
    p.font.color.rgb = MID_GREY
    p.font.name = ds.font_family
    return txBox


# ---------------------------------------------------------------------------
# Generic text
# ---------------------------------------------------------------------------
def add_textbox(slide, left, top, width, height, text, ds: DesignSystem,
                font_size=None, color=None, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, word_wrap=True):
    """Add a textbox with standard formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    if anchor:
        tf.paragraphs[0].alignment = alignment

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size or ds.font_body
    p.font.color.rgb = color or ds.secondary
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = ds.font_family
    p.alignment = alignment

    # Set vertical anchor
    from pptx.oxml.ns import qn
    txBody = txBox._element.txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None and anchor:
        anchor_map = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}
        bodyPr.set("anchor", anchor_map.get(anchor, "t"))

    return txBox


def add_bullet_list(slide, left, top, width, height, items, ds: DesignSystem,
                    font_size=None, color=None, spacing_pt=8, bullet_char="\u2022"):
    """Add a bulleted list with proper indent and spacing."""
    from lxml import etree
    from pptx.oxml.ns import qn as _qn

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    size = font_size or ds.font_body

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = color or ds.secondary
        p.font.name = ds.font_family
        p.space_after = Pt(spacing_pt)
        p.space_before = Pt(2)
        p.level = 0

        # Set indent and bullet via pPr
        pPr = p._pPr
        if pPr is None:
            pPr = etree.SubElement(p._p, _qn("a:pPr"))
        # Hanging indent: bullet at margin, text indented
        pPr.set("marL", str(int(Inches(0.25))))   # left margin for text
        pPr.set("indent", str(int(Inches(-0.2))))  # bullet hangs left of text

        # Bullet character
        buChar = etree.SubElement(pPr, _qn("a:buChar"))
        buChar.set("char", bullet_char)

        # Bullet color matches accent
        buClr = etree.SubElement(pPr, _qn("a:buClr"))
        srgb = etree.SubElement(buClr, _qn("a:srgbClr"))
        srgb.set("val", str(ds.accent1))

    return txBox


# ---------------------------------------------------------------------------
# Icons
# ---------------------------------------------------------------------------
def add_icon(slide, left, top, size, icon_name, ds: DesignSystem,
             color=None, bg_color=None, circle_bg=True):
    """
    Add an icon to a slide. Tries SVG file first, falls back to Unicode character.

    icon_name: name like "trending-up" (checks assets/icons/{name}.svg)
               or key from ICON_MAP, or a raw Unicode character.
    size: Inches value for the icon container (square).
    """
    import os
    from scripts.design_system import ICON_MAP

    icon_color = color or ds.accent1
    background = bg_color or LIGHT_GREY

    if circle_bg:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = background
        circle.line.fill.background()

    # Try to find SVG file in assets/icons/
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    svg_path = os.path.join(base_dir, "assets", "icons", f"{icon_name}.svg")

    if os.path.exists(svg_path):
        # Try to convert SVG to PNG and embed as image
        try:
            return _add_svg_icon(slide, left, top, size, svg_path, icon_color)
        except Exception:
            pass  # Fall through to Unicode

    # Unicode character fallback
    if len(icon_name) > 1:
        char = ICON_MAP.get(icon_name, icon_name)
    else:
        char = icon_name

    txBox = slide.shapes.add_textbox(left, top, size, size)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = char
    p.font.size = Pt(int(size / 914400 * 72 * 0.45))  # ~45% of container
    p.font.color.rgb = icon_color
    p.font.name = "Segoe UI Symbol"
    p.alignment = PP_ALIGN.CENTER

    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "ctr")

    return txBox


def _add_svg_icon(slide, left, top, size, svg_path, color):
    """Attempt to add an SVG icon as an embedded image. Requires cairosvg or Pillow."""
    import io
    try:
        import cairosvg
        # Render SVG to PNG at appropriate resolution
        png_size = int(size / 914400 * 96)  # Convert EMU to pixels at 96 DPI
        png_data = cairosvg.svg2png(
            url=svg_path,
            output_width=png_size,
            output_height=png_size,
        )
        # Pad: center the icon within the circle (80% of container)
        icon_size = int(size * 0.7)
        offset = int((size - icon_size) / 2)
        pic = slide.shapes.add_picture(
            io.BytesIO(png_data),
            left + offset, top + offset,
            icon_size, icon_size,
        )
        return pic
    except ImportError:
        raise  # Let caller fall back to Unicode


# ---------------------------------------------------------------------------
# Shapes
# ---------------------------------------------------------------------------
def add_rectangle(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rectangle(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------
def add_table(slide, left, top, width, height, headers, rows, ds: DesignSystem,
              col_widths=None):
    """
    Add a clean MBB-style table.
    headers: list of column header strings
    rows: list of lists (each inner list = one row of cell values)
    col_widths: optional list of Inches for each column
    """
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        col_w = int(width / num_cols)
        for i in range(num_cols):
            table.columns[i].width = col_w

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(header)
        _format_cell(cell, ds, bold=True, color=TABLE_HEADER_FG, fill=TABLE_HEADER_BG)

    # Data rows
    for row_idx, row_data in enumerate(rows):
        fill = TABLE_ALT_ROW if row_idx % 2 == 0 else WHITE
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            align = PP_ALIGN.RIGHT if is_numeric(value) else PP_ALIGN.LEFT
            _format_cell(cell, ds, fill=fill, alignment=align)

    # Apply MBB borders
    set_table_mbb_borders(table)

    return table_shape


def _format_cell(cell, ds, bold=False, color=None, fill=None, alignment=PP_ALIGN.LEFT):
    """Format a single table cell."""
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = ds.font_body
        paragraph.font.name = ds.font_family
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color or ds.secondary
        paragraph.alignment = alignment
    cell.text_frame.paragraphs[0].space_before = Pt(4)
    cell.text_frame.paragraphs[0].space_after = Pt(4)
    # Vertical centering
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        set_cell_fill(cell, fill)


# ---------------------------------------------------------------------------
# Charts
# ---------------------------------------------------------------------------
def add_bar_chart(slide, left, top, width, height, categories, series_list, ds: DesignSystem,
                  orientation="vertical", show_values=True, value_format=None):
    """
    Add a native bar/column chart.
    categories: list of category labels
    series_list: list of dicts with "name", "values", optional "color"
    """
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED if orientation == "vertical" else XL_CHART_TYPE.BAR_CLUSTERED
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s["name"], s["values"])

    chart_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = chart_frame.chart

    # Style the chart
    chart.has_legend = len(series_list) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = FONT_SMALL
        chart.legend.font.name = ds.font_family

    # Remove gridlines
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    chart.value_axis.visible = False

    # Style category axis
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = FONT_SMALL
    chart.category_axis.tick_labels.font.name = ds.font_family
    chart.category_axis.tick_labels.font.color.rgb = SLATE

    # Color series
    for i, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        if i < len(series_list) and "color" in series_list[i]:
            fill.fore_color.rgb = resolve_color(series_list[i]["color"])
        elif i < len(ds.viz_palette):
            fill.fore_color.rgb = ds.viz_palette[i]

        # Data labels
        if show_values:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.font.size = FONT_SMALL
            data_labels.font.name = ds.font_family
            data_labels.font.color.rgb = SLATE
            if value_format:
                data_labels.number_format = value_format
            if orientation == "vertical":
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            else:
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    return chart_frame


def add_line_chart(slide, left, top, width, height, categories, series_list, ds: DesignSystem,
                   show_values=False, value_format=None):
    """
    Add a native line chart.
    categories: list of category labels (e.g. time periods)
    series_list: list of dicts with "name", "values", optional "color"
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s["name"], s["values"])

    chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data)
    chart = chart_frame.chart

    chart.has_legend = len(series_list) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = FONT_SMALL
        chart.legend.font.name = ds.font_family

    # Light gridlines
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
    chart.value_axis.tick_labels.font.size = FONT_SMALL
    chart.value_axis.tick_labels.font.name = ds.font_family

    chart.category_axis.tick_labels.font.size = FONT_SMALL
    chart.category_axis.tick_labels.font.name = ds.font_family
    chart.category_axis.tick_labels.font.color.rgb = SLATE

    # Color series
    for i, series in enumerate(chart.series):
        line = series.format.line
        if i < len(series_list) and "color" in series_list[i]:
            line.color.rgb = resolve_color(series_list[i]["color"])
        elif i < len(ds.viz_palette):
            line.color.rgb = ds.viz_palette[i]
        line.width = Pt(2.5)

        if show_values:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.font.size = FONT_SMALL
            data_labels.font.name = ds.font_family
            if value_format:
                data_labels.number_format = value_format

    return chart_frame
