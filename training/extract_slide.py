#!/usr/bin/env python3
"""
Extract content from an existing MBB PowerPoint slide for use in the training loop.

Takes a .pptx file and extracts structured content (text, tables, charts)
into a JSON format compatible with plan.json.

Usage:
    python3 training/extract_slide.py --input reference.pptx --output extracted.json
    python3 training/extract_slide.py --input reference.pptx --slide 3  # Extract specific slide
"""

import argparse
import json
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_presentation(pptx_path, slide_number=None):
    """Extract content from a PPTX file into plan.json format."""
    prs = Presentation(pptx_path)

    result = {
        "title": os.path.splitext(os.path.basename(pptx_path))[0],
        "source_file": pptx_path,
        "slide_width_inches": prs.slide_width / 914400,
        "slide_height_inches": prs.slide_height / 914400,
        "slides": [],
    }

    slides_to_process = enumerate(prs.slides, 1)
    if slide_number:
        slides_to_process = [(slide_number, list(prs.slides)[slide_number - 1])]

    for idx, slide in slides_to_process:
        slide_data = extract_slide(slide, idx)
        result["slides"].append(slide_data)

    return result


def extract_slide(slide, slide_num):
    """Extract content from a single slide."""
    slide_data = {
        "slide_number": slide_num,
        "type": "unknown",  # Will be classified later
        "shapes": [],
        "text_content": [],
        "tables": [],
        "charts": [],
        "images": [],
        "notes": "",
    }

    # Extract speaker notes
    if slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        if notes_tf:
            slide_data["notes"] = notes_tf.text

    for shape in slide.shapes:
        shape_info = {
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "left_inches": shape.left / 914400 if shape.left else 0,
            "top_inches": shape.top / 914400 if shape.top else 0,
            "width_inches": shape.width / 914400 if shape.width else 0,
            "height_inches": shape.height / 914400 if shape.height else 0,
        }

        # Text content
        if shape.has_text_frame:
            paragraphs = []
            for para in shape.text_frame.paragraphs:
                para_info = {
                    "text": para.text,
                    "level": para.level,
                    "alignment": str(para.alignment) if para.alignment else None,
                    "font_size_pt": para.font.size / 12700 if para.font.size else None,
                    "font_bold": para.font.bold,
                    "font_italic": para.font.italic,
                    "font_name": para.font.name,
                    "font_color": str(para.font.color.rgb) if para.font.color and para.font.color.rgb else None,
                }
                paragraphs.append(para_info)

            shape_info["paragraphs"] = paragraphs
            slide_data["text_content"].append({
                "text": shape.text_frame.text,
                "position": {
                    "top": shape_info["top_inches"],
                    "left": shape_info["left_inches"],
                    "width": shape_info["width_inches"],
                    "height": shape_info["height_inches"],
                },
                "paragraphs": paragraphs,
            })

        # Tables
        if shape.has_table:
            table = shape.table
            headers = []
            rows = []
            for row_idx, row in enumerate(table.rows):
                row_data = [cell.text for cell in row.cells]
                if row_idx == 0:
                    headers = row_data
                else:
                    rows.append(row_data)
            slide_data["tables"].append({
                "headers": headers,
                "rows": rows,
                "position": {
                    "top": shape_info["top_inches"],
                    "left": shape_info["left_inches"],
                    "width": shape_info["width_inches"],
                    "height": shape_info["height_inches"],
                },
            })

        # Charts
        if shape.has_chart:
            chart = shape.chart
            chart_info = {
                "chart_type": str(chart.chart_type),
                "has_legend": chart.has_legend,
                "position": {
                    "top": shape_info["top_inches"],
                    "left": shape_info["left_inches"],
                    "width": shape_info["width_inches"],
                    "height": shape_info["height_inches"],
                },
                "series": [],
                "categories": [],
            }

            # Extract categories
            try:
                plot = chart.plots[0]
                cats = plot.categories
                if cats:
                    chart_info["categories"] = list(cats)
            except Exception:
                pass

            # Extract series data
            try:
                for series in chart.series:
                    series_info = {
                        "name": str(series.name) if series.name else "",
                        "values": list(series.values) if series.values else [],
                    }
                    chart_info["series"].append(series_info)
            except Exception:
                pass

            slide_data["charts"].append(chart_info)

        # Images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            slide_data["images"].append({
                "name": shape.name,
                "position": {
                    "top": shape_info["top_inches"],
                    "left": shape_info["left_inches"],
                    "width": shape_info["width_inches"],
                    "height": shape_info["height_inches"],
                },
            })

        slide_data["shapes"].append(shape_info)

    # Classify slide type
    slide_data["type"] = _classify_slide(slide_data)

    return slide_data


def _classify_slide(slide_data):
    """Heuristic classification of slide type based on content."""
    texts = slide_data["text_content"]
    tables = slide_data["tables"]
    charts = slide_data["charts"]
    num_shapes = len(slide_data["shapes"])

    # Check for charts
    if charts:
        chart_type = charts[0].get("chart_type", "")
        if "BAR" in chart_type or "COLUMN" in chart_type:
            return "bar_chart"
        if "LINE" in chart_type:
            return "line_chart"
        return "chart_unknown"

    # Check for tables
    if tables:
        return "data_table"

    # Text-only classification
    if num_shapes <= 3:
        # Minimal shapes — could be title, section divider, or Q&A
        for t in texts:
            font_sizes = [p.get("font_size_pt") for p in t.get("paragraphs", []) if p.get("font_size_pt")]
            if font_sizes and max(font_sizes) >= 36:
                if t["position"]["top"] < 2:
                    return "title_hero"
                return "section_divider"
        return "qa_page"

    # Multiple text blocks
    if len(texts) >= 5:
        # Check for bullet lists
        has_bullets = any(
            any(p.get("level", 0) > 0 for p in t.get("paragraphs", []))
            for t in texts
        )
        if has_bullets and num_shapes <= 8:
            return "executive_summary"

    # Check for stat-heavy slides (large numbers)
    large_font_count = sum(
        1 for t in texts
        for p in t.get("paragraphs", [])
        if p.get("font_size_pt") and p["font_size_pt"] >= 36
    )
    if large_font_count >= 2:
        return "key_stat"

    return "unknown"


def to_plan_json(extracted):
    """
    Convert extracted content into a plan.json compatible format.
    This is a best-effort conversion that will need manual review.
    """
    plan = {
        "title": extracted["title"],
        "slides": [],
    }

    for slide in extracted["slides"]:
        slide_plan = _convert_slide_to_plan(slide)
        plan["slides"].append(slide_plan)

    return plan


def _convert_slide_to_plan(slide):
    """Convert a single extracted slide to plan.json format."""
    slide_type = slide["type"]

    # Find the headline (topmost large text)
    headline = ""
    texts_sorted = sorted(slide["text_content"], key=lambda t: t["position"]["top"])
    for t in texts_sorted:
        for p in t.get("paragraphs", []):
            if p.get("font_bold") and p.get("text") and len(p["text"]) > 10:
                headline = p["text"]
                break
        if headline:
            break

    plan = {
        "type": slide_type,
        "headline": headline,
        "_raw_extraction": slide,  # Keep raw data for manual review
    }

    # Type-specific conversion
    if slide_type == "data_table" and slide["tables"]:
        table = slide["tables"][0]
        plan["headers"] = table["headers"]
        plan["rows"] = table["rows"]

    elif slide_type in ("bar_chart", "line_chart") and slide["charts"]:
        chart = slide["charts"][0]
        plan["categories"] = chart["categories"]
        plan["series"] = [
            {"name": s["name"], "values": s["values"]}
            for s in chart["series"]
        ]

    elif slide_type == "executive_summary":
        bullets = []
        for t in texts_sorted[1:]:  # Skip headline
            for p in t.get("paragraphs", []):
                if p.get("text") and len(p["text"]) > 20:
                    bullets.append(p["text"])
        plan["bullets"] = bullets

    elif slide_type == "key_stat":
        stats = []
        for t in texts_sorted:
            for p in t.get("paragraphs", []):
                if p.get("font_size_pt") and p["font_size_pt"] >= 36 and p.get("text"):
                    stats.append({"value": p["text"], "label": ""})
        plan["stats"] = stats

    return plan


def main():
    parser = argparse.ArgumentParser(description="Extract content from MBB PowerPoint slides")
    parser.add_argument("--input", required=True, help="Input .pptx file")
    parser.add_argument("--output", help="Output .json file (default: stdout)")
    parser.add_argument("--slide", type=int, help="Extract specific slide number only")
    parser.add_argument("--plan", action="store_true", help="Convert to plan.json format")
    args = parser.parse_args()

    extracted = extract_presentation(args.input, args.slide)

    if args.plan:
        result = to_plan_json(extracted)
    else:
        result = extracted

    output_json = json.dumps(result, indent=2, default=str)

    if args.output:
        with open(args.output, "w") as f:
            f.write(output_json)
        print(f"Extracted to: {args.output}")
    else:
        print(output_json)


if __name__ == "__main__":
    main()
