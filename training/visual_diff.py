#!/usr/bin/env python3
"""
Visual diff between reference MBB slides and generated output.

Converts both PPTX files to images (using python-pptx thumbnail extraction
or LibreOffice conversion), then compares them pixel-by-pixel and generates
a diff report.

Usage:
    python3 training/visual_diff.py --reference ref.pptx --generated gen.pptx --output diffs/
    python3 training/visual_diff.py --reference ref.pptx --generated gen.pptx --slide 3

Requirements:
    pip install Pillow
    LibreOffice (for PPTX → PNG conversion): brew install --cask libreoffice
"""

import argparse
import json
import os
import subprocess
import sys
import tempfile
import shutil
import glob

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))


def pptx_to_images(pptx_path, output_dir):
    """
    Convert PPTX to PNG images using LibreOffice.
    Returns list of image paths sorted by slide number.
    """
    os.makedirs(output_dir, exist_ok=True)

    # Try LibreOffice first
    lo_paths = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "soffice",
        "libreoffice",
    ]

    lo_binary = None
    for path in lo_paths:
        if shutil.which(path) or os.path.exists(path):
            lo_binary = path
            break

    if lo_binary is None:
        print("WARNING: LibreOffice not found. Install with: brew install --cask libreoffice")
        print("Falling back to metadata-only comparison.")
        return []

    # Convert to PDF first, then to images
    with tempfile.TemporaryDirectory() as tmpdir:
        # PPTX → PDF
        result = subprocess.run(
            [lo_binary, "--headless", "--convert-to", "png", "--outdir", output_dir, pptx_path],
            capture_output=True, text=True, timeout=120,
        )

        if result.returncode != 0:
            # Try PDF route
            subprocess.run(
                [lo_binary, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path],
                capture_output=True, text=True, timeout=120,
            )

            pdf_files = glob.glob(os.path.join(tmpdir, "*.pdf"))
            if pdf_files:
                # PDF → PNG (requires Pillow or another tool)
                try:
                    from pdf2image import convert_from_path
                    images = convert_from_path(pdf_files[0], dpi=150)
                    for i, img in enumerate(images):
                        img.save(os.path.join(output_dir, f"slide_{i+1:02d}.png"))
                except ImportError:
                    print("WARNING: pdf2image not available. pip install pdf2image")
                    return []

    # Collect output images
    images = sorted(glob.glob(os.path.join(output_dir, "*.png")))
    return images


def compare_images(ref_path, gen_path, diff_path):
    """
    Compare two images and generate a visual diff.
    Returns similarity score (0.0 = identical, 1.0 = completely different).
    """
    try:
        from PIL import Image, ImageChops, ImageDraw, ImageFont
    except ImportError:
        print("WARNING: Pillow not available. pip install Pillow")
        return None

    ref_img = Image.open(ref_path).convert("RGB")
    gen_img = Image.open(gen_path).convert("RGB")

    # Resize to same dimensions if needed
    if ref_img.size != gen_img.size:
        gen_img = gen_img.resize(ref_img.size, Image.LANCZOS)

    # Pixel diff
    diff = ImageChops.difference(ref_img, gen_img)

    # Calculate similarity score
    pixels = list(diff.getdata())
    total_diff = sum(sum(p) for p in pixels)
    max_diff = len(pixels) * 3 * 255  # 3 channels, max 255 per channel
    similarity = total_diff / max_diff if max_diff > 0 else 0

    # Create side-by-side comparison image
    width = ref_img.width
    height = ref_img.height
    comparison = Image.new("RGB", (width * 3, height + 40), (255, 255, 255))

    # Headers
    draw = ImageDraw.Draw(comparison)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
    except (OSError, IOError):
        font = ImageFont.load_default()

    draw.text((width * 0 + 10, 10), "REFERENCE", fill=(0, 0, 0), font=font)
    draw.text((width * 1 + 10, 10), "GENERATED", fill=(0, 0, 0), font=font)
    draw.text((width * 2 + 10, 10), f"DIFF (score: {similarity:.4f})", fill=(220, 38, 38), font=font)

    comparison.paste(ref_img, (0, 40))
    comparison.paste(gen_img, (width, 40))

    # Amplify diff for visibility
    amplified_diff = diff.point(lambda x: min(x * 5, 255))
    comparison.paste(amplified_diff, (width * 2, 40))

    comparison.save(diff_path)

    return similarity


def compare_metadata(ref_pptx, gen_pptx):
    """
    Compare structural metadata between two PPTX files.
    Useful when image comparison is not available.
    """
    from pptx import Presentation

    ref = Presentation(ref_pptx)
    gen = Presentation(gen_pptx)

    report = {
        "ref_slides": len(ref.slides),
        "gen_slides": len(gen.slides),
        "slide_count_match": len(ref.slides) == len(gen.slides),
        "slides": [],
    }

    max_slides = max(len(ref.slides), len(gen.slides))
    ref_slides = list(ref.slides)
    gen_slides = list(gen.slides)

    for i in range(max_slides):
        slide_report = {"slide_number": i + 1}

        if i < len(ref_slides):
            ref_slide = ref_slides[i]
            slide_report["ref_shapes"] = len(ref_slide.shapes)
            slide_report["ref_text"] = _get_slide_text(ref_slide)
        else:
            slide_report["ref_shapes"] = 0
            slide_report["ref_text"] = ""

        if i < len(gen_slides):
            gen_slide = gen_slides[i]
            slide_report["gen_shapes"] = len(gen_slide.shapes)
            slide_report["gen_text"] = _get_slide_text(gen_slide)
        else:
            slide_report["gen_shapes"] = 0
            slide_report["gen_text"] = ""

        # Text similarity (simple Jaccard)
        ref_words = set(slide_report["ref_text"].lower().split())
        gen_words = set(slide_report["gen_text"].lower().split())
        if ref_words or gen_words:
            intersection = ref_words & gen_words
            union = ref_words | gen_words
            slide_report["text_similarity"] = len(intersection) / len(union) if union else 1.0
        else:
            slide_report["text_similarity"] = 1.0

        report["slides"].append(slide_report)

    return report


def _get_slide_text(slide):
    """Extract all text from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return " ".join(texts)


def run_diff(ref_path, gen_path, output_dir, slide_number=None):
    """Run full visual diff pipeline."""
    os.makedirs(output_dir, exist_ok=True)

    print(f"Reference: {ref_path}")
    print(f"Generated: {gen_path}")
    print(f"Output:    {output_dir}")
    print()

    # Metadata comparison (always available)
    meta_report = compare_metadata(ref_path, gen_path)
    meta_path = os.path.join(output_dir, "metadata_report.json")
    with open(meta_path, "w") as f:
        json.dump(meta_report, f, indent=2)
    print(f"Metadata report: {meta_path}")

    # Print text similarity summary
    print(f"\nSlide count: ref={meta_report['ref_slides']}, gen={meta_report['gen_slides']}")
    for s in meta_report["slides"]:
        sim = s.get("text_similarity", 0)
        status = "OK" if sim > 0.7 else "LOW" if sim > 0.3 else "MISMATCH"
        print(f"  Slide {s['slide_number']}: text similarity={sim:.2f} [{status}] "
              f"(ref shapes: {s['ref_shapes']}, gen shapes: {s['gen_shapes']})")

    # Visual comparison (requires LibreOffice + Pillow)
    print("\nAttempting visual comparison...")
    ref_img_dir = os.path.join(output_dir, "ref_images")
    gen_img_dir = os.path.join(output_dir, "gen_images")

    ref_images = pptx_to_images(ref_path, ref_img_dir)
    gen_images = pptx_to_images(gen_path, gen_img_dir)

    if ref_images and gen_images:
        print(f"  Reference images: {len(ref_images)}")
        print(f"  Generated images: {len(gen_images)}")

        visual_scores = []
        for i in range(min(len(ref_images), len(gen_images))):
            if slide_number and (i + 1) != slide_number:
                continue
            diff_path = os.path.join(output_dir, f"diff_slide_{i+1:02d}.png")
            score = compare_images(ref_images[i], gen_images[i], diff_path)
            if score is not None:
                visual_scores.append({"slide": i + 1, "score": score, "diff_image": diff_path})
                status = "GOOD" if score < 0.05 else "FAIR" if score < 0.15 else "POOR"
                print(f"  Slide {i+1}: visual diff score={score:.4f} [{status}] → {diff_path}")

        # Save visual report
        visual_path = os.path.join(output_dir, "visual_report.json")
        with open(visual_path, "w") as f:
            json.dump(visual_scores, f, indent=2)
    else:
        print("  Visual comparison skipped (LibreOffice not available)")

    # Overall summary
    print(f"\nDiff results saved to: {output_dir}")
    return meta_report


def main():
    parser = argparse.ArgumentParser(description="Visual diff between reference and generated PPTX")
    parser.add_argument("--reference", required=True, help="Reference PPTX file")
    parser.add_argument("--generated", required=True, help="Generated PPTX file")
    parser.add_argument("--output", default="training/diffs", help="Output directory for diff results")
    parser.add_argument("--slide", type=int, help="Compare specific slide only")
    args = parser.parse_args()

    run_diff(args.reference, args.generated, args.output, args.slide)


if __name__ == "__main__":
    main()
